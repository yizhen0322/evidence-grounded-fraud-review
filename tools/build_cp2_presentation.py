"""Build the final rubric-aligned CP2 viva presentation with speaker notes."""

from __future__ import annotations

import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
CONTENT = ROOT / "reports/presentation/CP2_Presentation_Content.md"
OUTPUT = ROOT / "reports/presentation/CP2_Presentation_Ng_Yi_Zhen.pptx"
SCRIPT_OUTPUT = ROOT / "reports/presentation/CP2_Final_Presentation_Script.md"

PAPER = RGBColor(247, 248, 246)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(18, 34, 55)
BLUE = RGBColor(45, 94, 160)
BLUE_SOFT = RGBColor(223, 232, 243)
TEAL = RGBColor(29, 122, 113)
TEAL_SOFT = RGBColor(222, 239, 235)
AMBER = RGBColor(188, 115, 44)
AMBER_SOFT = RGBColor(246, 234, 218)
RUST = RGBColor(164, 70, 58)
RUST_SOFT = RGBColor(244, 226, 222)
INK = RGBColor(35, 43, 52)
MUTED = RGBColor(91, 102, 113)
LINE = RGBColor(211, 217, 220)
PALE = RGBColor(236, 239, 239)

FIGURES = ROOT / "reports/figures"
WORKBENCH = FIGURES / "workbench_investigation.png"
GUARDRAIL_FAILURE = FIGURES / "workbench_guardrail_failure.png"
DETECTOR_BARS = FIGURES / "detector_metric_bars.png"
NARRATIVE_BARS = FIGURES / "narrative_delivery_bars.png"
HUMAN_PREFERENCES = FIGURES / "human_eval_preferences.png"


def parse_content() -> list[dict[str, object]]:
    text = CONTENT.read_text(encoding="utf-8")
    blocks = re.split(r"(?=^## Slide \d+\.)", text, flags=re.MULTILINE)
    slides: list[dict[str, object]] = []
    for block in blocks:
        match = re.match(r"## Slide (\d+)\. (.+)", block)
        if not match:
            continue
        number = int(match.group(1))
        title = match.group(2).strip()
        bullet_region = block.split("**Existing visual recommendation:**", 1)[0]
        bullets = re.findall(r"^- (.+)$", bullet_region, flags=re.MULTILINE)
        script_match = re.search(
            r"\*\*Speaker script:\*\*\s*\n(.+?)(?=\n## Slide|\n# Examiner|\Z)",
            block,
            flags=re.DOTALL,
        )
        script = " ".join(script_match.group(1).strip().split()) if script_match else ""
        slides.append({"number": number, "title": title, "bullets": bullets, "script": script})
    if len(slides) != 12:
        raise ValueError(f"Expected 12 slides, found {len(slides)}")
    return slides


def clean(text: str) -> str:
    return text.replace("**", "").replace("`", "").replace("+/-", "±")


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text: str = "",
    *,
    size: float = 18,
    colour=INK,
    bold: bool = False,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 2,
):
    shape = slide.shapes.add_textbox(x, y, w, h)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Pt(margin)
    frame.margin_right = Pt(margin)
    frame.margin_top = Pt(margin)
    frame.margin_bottom = Pt(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = colour
    return shape


def add_rule(slide, x, y, w, colour=LINE, height=1.2) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Pt(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()


def add_header(slide, number: int, title: str, kicker: str) -> None:
    add_rule(slide, 0, 0, Inches(13.333), NAVY, height=5)
    add_textbox(slide, Inches(0.64), Inches(0.22), Inches(3.5), Inches(0.25), kicker.upper(),
                size=10, colour=TEAL, bold=True)
    add_textbox(slide, Inches(0.64), Inches(0.52), Inches(11.7), Inches(0.42), title,
                size=26, colour=NAVY, bold=True, font="Aptos Display")
    add_textbox(slide, Inches(12.35), Inches(0.28), Inches(0.38), Inches(0.24), f"{number:02d}",
                size=11, colour=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_footer(slide) -> None:
    add_rule(slide, Inches(0.64), Inches(7.17), Inches(12.05), LINE, height=0.9)
    add_textbox(slide, Inches(0.66), Inches(7.21), Inches(6.2), Inches(0.16),
                "Ng Yi Zhen | Capstone Project 2", size=8, colour=MUTED)


def add_panel(slide, x, y, w, h, *, fill=WHITE, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(0.9)
    return shape


def add_badge(slide, x, y, text: str, *, fill=NAVY, colour=WHITE, w=1.15) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(w), Inches(0.34))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_textbox(slide, x + Inches(0.05), y + Inches(0.055), Inches(w - 0.10), Inches(0.2), text,
                size=9, colour=colour, bold=True, align=PP_ALIGN.CENTER)


def add_claim_banner(slide, text: str, *, y=1.08, fill=NAVY, colour=WHITE, h=0.62) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.66), Inches(y), Inches(12.0), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    add_textbox(slide, Inches(0.95), Inches(y + 0.15), Inches(11.4), Inches(h - 0.2), text,
                size=17, colour=colour, bold=True, align=PP_ALIGN.CENTER)


def add_metric_card(slide, x, y, w, h, value: str, label: str, *, colour=BLUE, note: str = "") -> None:
    add_panel(slide, x, y, w, h, fill=WHITE, line=LINE)
    add_rule(slide, x, y, w, colour, height=5)
    add_textbox(slide, x + Inches(0.16), y + Inches(0.24), w - Inches(0.32), Inches(0.55), value,
                size=25, colour=colour, bold=True)
    add_textbox(slide, x + Inches(0.16), y + Inches(0.86), w - Inches(0.32), Inches(0.38), label,
                size=12, colour=INK, bold=True)
    if note:
        add_textbox(slide, x + Inches(0.16), y + Inches(1.27), w - Inches(0.32), h - Inches(1.38), note,
                    size=9.5, colour=MUTED)


def add_picture_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        iw, ih = image.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    return slide.shapes.add_picture(str(path), px, py, pw, ph)


def add_picture_crop(slide, path: Path, x, y, w, h, *, left=0.0, top=0.0, right=0.0, bottom=0.0):
    picture = slide.shapes.add_picture(str(path), x, y, w, h)
    picture.crop_left = left
    picture.crop_top = top
    picture.crop_right = right
    picture.crop_bottom = bottom
    return picture


def add_arrow(slide, x, y, w=0.45, text="→", colour=MUTED, size=22) -> None:
    add_textbox(slide, Inches(x), Inches(y), Inches(w), Inches(0.32), text,
                size=size, colour=colour, bold=True, align=PP_ALIGN.CENTER)


def add_check_row(slide, x, y, label: str, *, status="yes", width=2.2) -> None:
    colours = {"yes": (TEAL, TEAL_SOFT, "YES"), "partial": (AMBER, AMBER_SOFT, "PARTIAL"), "no": (MUTED, PALE, "NO")}
    colour, fill, text = colours[status]
    add_panel(slide, x, y, Inches(width), Inches(0.42), fill=fill, line=fill, radius=False)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.10), Inches(width - 0.75), Inches(0.2), label,
                size=9.5, colour=INK)
    add_textbox(slide, x + Inches(width - 0.62), y + Inches(0.10), Inches(0.52), Inches(0.2), text,
                size=8.5, colour=colour, bold=True, align=PP_ALIGN.RIGHT)


def add_notes(slide, script: str) -> None:
    slide.notes_slide.notes_text_frame.text = script


def build_title_slide(slide, script: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_textbox(slide, Inches(0.72), Inches(0.58), Inches(4.8), Inches(0.3),
                "SUNWAY UNIVERSITY | CAPSTONE PROJECT 2", size=11, colour=RGBColor(164, 197, 211), bold=True)
    add_textbox(slide, Inches(0.72), Inches(1.28), Inches(5.3), Inches(1.65),
                "Fraud alert\nexplanations that\nfail closed", size=31, colour=WHITE, bold=True,
                font="Aptos Display")
    add_textbox(slide, Inches(0.74), Inches(3.34), Inches(4.95), Inches(0.82),
                "Evidence-grounded local-LLM delivery for credit-card alert review",
                size=17, colour=RGBColor(211, 223, 230))
    add_panel(slide, Inches(0.72), Inches(4.48), Inches(4.9), Inches(1.08), fill=RGBColor(27, 49, 74), line=RGBColor(55, 81, 105))
    add_textbox(slide, Inches(0.98), Inches(4.75), Inches(4.38), Inches(0.52),
                "The LLM writes only when deterministic checks confirm that it stayed inside the stored evidence.",
                size=14, colour=WHITE, bold=True)
    add_textbox(slide, Inches(0.74), Inches(6.45), Inches(4.5), Inches(0.45),
                "Ng Yi Zhen | 23076003", size=16, colour=WHITE, bold=True)

    add_panel(slide, Inches(6.15), Inches(0.58), Inches(6.55), Inches(6.35), fill=WHITE, line=RGBColor(56, 76, 96))
    add_picture_crop(slide, WORKBENCH, Inches(6.30), Inches(0.73), Inches(6.25), Inches(5.95),
                     left=0.14, top=0.04, right=0.01, bottom=0.05)
    add_badge(slide, Inches(10.88), Inches(6.27), "LOCAL REVIEW", fill=TEAL, w=1.45)
    add_notes(slide, script)


def build_problem_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Problem")
    add_claim_banner(slide, "Readable does not mean faithful.", y=1.12, fill=NAVY)
    columns = [
        ("1", "Rare-event detector", "Accuracy can look strong while the model misses fraud. Use AP, precision, recall, F1, and confusion counts.", BLUE, BLUE_SOFT),
        ("2", "Technical evidence", "SHAP gives signed attribution, but raw feature contributions are repetitive and difficult to review quickly.", TEAL, TEAL_SOFT),
        ("3", "Narrative risk", "A readable LLM summary may omit evidence, add an unsupported feature, or reverse a contribution direction.", AMBER, AMBER_SOFT),
    ]
    for index, (n, heading, body, colour, soft) in enumerate(columns):
        x = Inches(0.72 + index * 4.13)
        add_panel(slide, x, Inches(2.12), Inches(3.72), Inches(3.58), fill=WHITE, line=LINE)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.22), Inches(2.36), Inches(0.62), Inches(0.62))
        circle.fill.solid(); circle.fill.fore_color.rgb = soft; circle.line.fill.background()
        add_textbox(slide, x + Inches(0.22), Inches(2.51), Inches(0.62), Inches(0.22), n,
                    size=13, colour=colour, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Inches(0.22), Inches(3.18), Inches(3.20), Inches(0.42), heading,
                    size=17, colour=colour, bold=True)
        add_textbox(slide, x + Inches(0.22), Inches(3.78), Inches(3.20), Inches(1.35), body,
                    size=13, colour=INK)
    add_panel(slide, Inches(1.85), Inches(6.05), Inches(9.65), Inches(0.68), fill=PALE, line=LINE, radius=False)
    add_textbox(slide, Inches(2.15), Inches(6.25), Inches(9.05), Inches(0.26),
                "Research problem: keep generated explanations inside a verifiable model-evidence contract.",
                size=15, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_gap_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Literature and research questions")
    add_badge(slide, Inches(0.66), Inches(1.14), "REVIEWED LITERATURE", fill=BLUE, w=1.75)
    headers = ["Study", "Code validator", "Fail-closed fallback", "Local model", "Raw vs delivered"]
    rows = [
        ("AlMarri et al. (2025)", "Not reported", "Not reported", "Yes", "Not reported"),
        ("Zytek et al. (2024)", "No", "Partial", "Partial", "Not reported"),
        ("Bello et al. (2025)", "Not reported", "Not reported", "Partial", "Not reported"),
        ("Martens et al. (2025)", "Not reported", "Not reported", "No", "Not reported"),
        ("This project", "Yes", "Yes", "Yes", "Yes"),
    ]
    x0, y0 = Inches(0.66), Inches(1.62)
    widths = [Inches(2.45), Inches(2.25), Inches(2.55), Inches(2.15), Inches(2.60)]
    row_h = Inches(0.48)
    for col, header in enumerate(headers):
        x = x0 + sum(widths[:col])
        box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y0, widths[col], row_h)
        box.fill.solid(); box.fill.fore_color.rgb = NAVY; box.line.color.rgb = WHITE
        add_textbox(slide, x + Inches(0.05), y0 + Inches(0.12), widths[col] - Inches(0.10), Inches(0.20), header,
                    size=9.5, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for row_index, row in enumerate(rows, start=1):
        for col, value in enumerate(row):
            x = x0 + sum(widths[:col])
            y = y0 + row_h * row_index
            is_project = row_index == len(rows)
            fill = TEAL_SOFT if is_project else (WHITE if row_index % 2 else PALE)
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, widths[col], row_h)
            box.fill.solid(); box.fill.fore_color.rgb = fill; box.line.color.rgb = LINE
            colour = TEAL if value == "Yes" else (AMBER if value == "Partial" else MUTED)
            add_textbox(slide, x + Inches(0.05), y + Inches(0.12), widths[col] - Inches(0.10), Inches(0.20), value,
                        size=9, colour=colour if col else INK, bold=is_project or value in {"Yes", "Partial"},
                        align=PP_ALIGN.CENTER if col else PP_ALIGN.LEFT)
    add_textbox(slide, Inches(0.72), Inches(4.63), Inches(11.8), Inches(0.28),
                "Within the reviewed literature, none of the compared studies evaluated all four properties together.",
                size=11, colour=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(0.72), Inches(5.12), Inches(5.88), Inches(1.45), fill=BLUE_SOFT, line=BLUE_SOFT)
    add_textbox(slide, Inches(0.98), Inches(5.36), Inches(0.75), Inches(0.26), "RQ1", size=13, colour=BLUE, bold=True)
    add_textbox(slide, Inches(1.75), Inches(5.30), Inches(4.50), Inches(0.70),
                "How effectively do deterministic checks detect violations and enforce fallback?",
                size=14, colour=INK, bold=True)
    add_panel(slide, Inches(6.74), Inches(5.12), Inches(5.88), Inches(1.45), fill=TEAL_SOFT, line=TEAL_SOFT)
    add_textbox(slide, Inches(7.00), Inches(5.36), Inches(0.75), Inches(0.26), "RQ2", size=13, colour=TEAL, bold=True)
    add_textbox(slide, Inches(7.77), Inches(5.30), Inches(4.48), Inches(0.70),
                "What does guarded local-LLM delivery add, preserve, or lose?",
                size=14, colour=INK, bold=True)
    add_footer(slide)
    add_notes(slide, script)


def build_roles_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Study design")
    add_panel(slide, Inches(0.72), Inches(1.35), Inches(5.05), Inches(4.82), fill=WHITE, line=BLUE)
    add_badge(slide, Inches(1.00), Inches(1.62), "ULB | REAL DATA", fill=BLUE, w=1.55)
    add_textbox(slide, Inches(1.00), Inches(2.15), Inches(4.45), Inches(0.48), "Real anonymised benchmark",
                size=21, colour=BLUE, bold=True)
    add_textbox(slide, Inches(1.00), Inches(2.93), Inches(4.38), Inches(1.8),
                "• Detector comparison across six groups and five seeds\n\n• Leakage-controlled SHAP evidence chain\n\n• Anonymous-feature prompt stress test",
                size=14, colour=INK)
    add_textbox(slide, Inches(1.00), Inches(5.35), Inches(4.35), Inches(0.45),
                "Boundary: V1 to V28 are faithful labels, but not business-readable semantics.",
                size=11, colour=MUTED)

    add_panel(slide, Inches(7.56), Inches(1.35), Inches(5.05), Inches(4.82), fill=WHITE, line=TEAL)
    add_badge(slide, Inches(7.84), Inches(1.62), "S0 | SYNTHETIC", fill=TEAL, w=1.55)
    add_textbox(slide, Inches(7.84), Inches(2.15), Inches(4.45), Inches(0.48), "Readable semantic stream",
                size=21, colour=TEAL, bold=True)
    add_textbox(slide, Inches(7.84), Inches(2.93), Inches(4.38), Inches(1.8),
                "• Operationally meaningful transaction fields\n\n• Structured explanation and fallback evaluation\n\n• Workbench and human-pilot stimuli",
                size=14, colour=INK)
    add_textbox(slide, Inches(7.84), Inches(5.35), Inches(4.35), Inches(0.45),
                "Boundary: reproducible simulation, not evidence of real-bank deployment validity.",
                size=11, colour=MUTED)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.94), Inches(2.55), Inches(1.44), Inches(1.44))
    circle.fill.solid(); circle.fill.fore_color.rgb = NAVY; circle.line.fill.background()
    add_textbox(slide, Inches(6.05), Inches(2.91), Inches(1.22), Inches(0.50), "ONE\nBOUNDARY",
                size=12, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_arrow(slide, 5.63, 3.02, text="→", colour=BLUE)
    add_arrow(slide, 7.25, 3.02, text="→", colour=TEAL)
    add_panel(slide, Inches(2.15), Inches(6.38), Inches(9.05), Inches(0.52), fill=PALE, line=LINE, radius=False)
    add_textbox(slide, Inches(2.35), Inches(6.52), Inches(8.65), Inches(0.22),
                "The detector scores are not directly compared. Both contexts test whether generated text stays inside stored evidence.",
                size=12, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_method_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Methodology")
    add_claim_banner(slide, "The evidence chain is frozen before the language model is allowed to write.", y=1.08, fill=NAVY)
    stages = [
        ("01", "Prepare", "Deduplicate\nStable case IDs\nSplit first", BLUE),
        ("02", "Benchmark", "6 detector groups\n5 fixed seeds", BLUE),
        ("03", "Freeze", "Selected model\nThreshold\nPredictions", NAVY),
        ("04", "Explain", "Signed SHAP\nTop-3 reason codes", TEAL),
        ("05", "Generate", "Minimised payload\nRaw candidate kept", AMBER),
        ("06", "Deliver", "Validate unchanged\nor deterministic fallback", RUST),
    ]
    for index, (n, heading, body, colour) in enumerate(stages):
        x = 0.66 + index * 2.05
        add_panel(slide, Inches(x), Inches(2.10), Inches(1.78), Inches(2.25), fill=WHITE, line=colour)
        add_textbox(slide, Inches(x + 0.15), Inches(2.30), Inches(0.55), Inches(0.25), n,
                    size=10, colour=colour, bold=True)
        add_textbox(slide, Inches(x + 0.15), Inches(2.72), Inches(1.48), Inches(0.34), heading,
                    size=15, colour=colour, bold=True)
        add_textbox(slide, Inches(x + 0.15), Inches(3.28), Inches(1.48), Inches(0.80), body,
                    size=11.5, colour=INK)
        if index < len(stages) - 1:
            add_arrow(slide, x + 1.78, 3.03, w=0.27, size=18)
    add_panel(slide, Inches(0.72), Inches(4.82), Inches(8.15), Inches(1.42), fill=BLUE_SOFT, line=BLUE_SOFT)
    add_badge(slide, Inches(0.98), Inches(5.04), "IMMUTABLE RESEARCH EVIDENCE", fill=BLUE, w=2.20)
    add_textbox(slide, Inches(0.98), Inches(5.52), Inches(7.55), Inches(0.42),
                "Manifests, predictions, metrics, SHAP reason codes, raw candidates, validator results, and delivered briefs",
                size=12.5, colour=INK)
    add_panel(slide, Inches(9.06), Inches(4.82), Inches(3.58), Inches(1.42), fill=TEAL_SOFT, line=TEAL_SOFT)
    add_badge(slide, Inches(9.32), Inches(5.04), "WRITABLE WORKFLOW", fill=TEAL, w=1.75)
    add_textbox(slide, Inches(9.32), Inches(5.52), Inches(3.05), Inches(0.42),
                "Analyst status, notes, and activity only",
                size=12.5, colour=INK)
    add_textbox(slide, Inches(0.88), Inches(6.58), Inches(11.5), Inches(0.28),
                "Scaling, SMOTE, autoencoder fitting, threshold selection, and SHAP never use the test set during training decisions.",
                size=11, colour=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_detector_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "ULB supporting benchmark")
    add_badge(slide, Inches(0.68), Inches(1.12), "ULB", fill=BLUE, w=0.70)
    add_picture_contain(slide, DETECTOR_BARS, Inches(0.62), Inches(1.42), Inches(8.35), Inches(5.25))
    add_metric_card(slide, Inches(9.18), Inches(1.34), Inches(1.58), Inches(1.72), "0.855", "G6 mean AP", colour=BLUE)
    add_metric_card(slide, Inches(10.92), Inches(1.34), Inches(1.58), Inches(1.72), "0.870", "G2 mean F1", colour=TEAL)
    add_metric_card(slide, Inches(9.18), Inches(3.24), Inches(1.58), Inches(1.72), "0.817", "G7 recall", colour=AMBER)
    add_metric_card(slide, Inches(10.92), Inches(3.24), Inches(1.58), Inches(1.72), "2.8", "G6 mean FP", colour=RUST)
    add_panel(slide, Inches(9.18), Inches(5.20), Inches(3.32), Inches(1.20), fill=PALE, line=LINE, radius=False)
    add_textbox(slide, Inches(9.42), Inches(5.42), Inches(2.84), Inches(0.30), "Selection decision", size=12, colour=NAVY, bold=True)
    add_textbox(slide, Inches(9.42), Inches(5.82), Inches(2.84), Inches(0.40),
                "G6 supports fewer false alerts. G7 misses fewer fraud cases. Neither is universally better.",
                size=10.5, colour=INK)
    add_textbox(slide, Inches(0.92), Inches(6.65), Inches(11.5), Inches(0.25),
                "Autoencoder features did not show a clear general advantage. G6 was frozen as a reproducible evidence source.",
                size=12.5, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_innovation_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Primary contribution")
    add_claim_banner(slide, "Generation quality and delivery safety are measured separately.", y=1.08, fill=NAVY)
    add_panel(slide, Inches(0.72), Inches(2.02), Inches(3.45), Inches(3.80), fill=AMBER_SOFT, line=AMBER)
    add_badge(slide, Inches(0.98), Inches(2.28), "OFF POLICY", fill=AMBER, w=1.10)
    add_textbox(slide, Inches(0.98), Inches(2.84), Inches(2.92), Inches(0.38), "Retain the raw candidate", size=17, colour=AMBER, bold=True)
    add_panel(slide, Inches(0.98), Inches(3.45), Inches(2.92), Inches(1.18), fill=WHITE, line=LINE, radius=False)
    add_textbox(slide, Inches(1.15), Inches(3.65), Inches(2.55), Inches(0.72),
                "Generated text\nACTION section missing\nCandidate preserved for measurement",
                size=11.5, colour=INK)
    add_textbox(slide, Inches(0.98), Inches(4.94), Inches(2.92), Inches(0.45),
                "The experiment records what the model actually returned.",
                size=11.5, colour=MUTED)

    add_panel(slide, Inches(4.58), Inches(2.02), Inches(4.18), Inches(3.80), fill=WHITE, line=TEAL)
    add_badge(slide, Inches(4.86), Inches(2.28), "DETERMINISTIC GATE", fill=TEAL, w=1.75)
    add_textbox(slide, Inches(4.86), Inches(2.84), Inches(3.62), Inches(0.38), "Validate the same candidate", size=17, colour=TEAL, bold=True)
    add_check_row(slide, Inches(4.88), Inches(3.45), "Required structure", status="yes", width=3.55)
    add_check_row(slide, Inches(4.88), Inches(3.98), "Evidence completeness", status="yes", width=3.55)
    add_check_row(slide, Inches(4.88), Inches(4.51), "Feature grounding", status="yes", width=3.55)
    add_check_row(slide, Inches(4.88), Inches(5.04), "Contribution direction", status="yes", width=3.55)

    add_panel(slide, Inches(9.18), Inches(2.02), Inches(3.45), Inches(3.80), fill=BLUE_SOFT, line=BLUE)
    add_badge(slide, Inches(9.44), Inches(2.28), "ON POLICY", fill=BLUE, w=1.10)
    add_textbox(slide, Inches(9.44), Inches(2.84), Inches(2.92), Inches(0.38), "Control what is delivered", size=17, colour=BLUE, bold=True)
    add_panel(slide, Inches(9.44), Inches(3.45), Inches(2.92), Inches(0.82), fill=TEAL_SOFT, line=TEAL, radius=False)
    add_textbox(slide, Inches(9.62), Inches(3.69), Inches(2.56), Inches(0.24), "PASS  |  deliver unchanged", size=12, colour=TEAL, bold=True)
    add_panel(slide, Inches(9.44), Inches(4.45), Inches(2.92), Inches(0.82), fill=RUST_SOFT, line=RUST, radius=False)
    add_textbox(slide, Inches(9.62), Inches(4.69), Inches(2.56), Inches(0.24), "FAIL  |  evidence-derived fallback", size=12, colour=RUST, bold=True)
    add_arrow(slide, 4.17, 3.62, w=0.40, size=22)
    add_arrow(slide, 8.77, 3.62, w=0.40, size=22)
    add_panel(slide, Inches(1.65), Inches(6.16), Inches(10.02), Inches(0.58), fill=RUST_SOFT, line=RUST_SOFT, radius=False)
    add_textbox(slide, Inches(1.90), Inches(6.33), Inches(9.52), Inches(0.24),
                "A normaliser can hide a generation failure. This design rejects the candidate and rebuilds fallback from the source evidence.",
                size=12.5, colour=RUST, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_guardrail_results_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "ULB explanation experiment")
    add_metric_card(slide, Inches(0.70), Inches(1.18), Inches(3.75), Inches(1.60), "2 / 51", "Strict prompt raw violations", colour=BLUE,
                    note="3.92% detected-any violation")
    add_metric_card(slide, Inches(4.78), Inches(1.18), Inches(3.75), Inches(1.60), "51 / 51", "Simple prompt raw violations", colour=AMBER,
                    note="100% detected-any violation")
    add_metric_card(slide, Inches(8.86), Inches(1.18), Inches(3.75), Inches(1.60), "0 / 49", "Manual semantic audit", colour=TEAL,
                    note="95% Wilson CI 0% to 7.27%; one reviewer")
    add_picture_contain(slide, NARRATIVE_BARS, Inches(0.70), Inches(3.08), Inches(7.72), Inches(3.72))
    add_panel(slide, Inches(8.72), Inches(3.12), Inches(3.89), Inches(3.47), fill=WHITE, line=LINE)
    add_badge(slide, Inches(8.98), Inches(3.38), "CALIBRATION CORPUS", fill=NAVY, w=1.68)
    add_textbox(slide, Inches(8.98), Inches(4.04), Inches(1.32), Inches(0.45), "330/330", size=24, colour=RUST, bold=True)
    add_textbox(slide, Inches(10.34), Inches(4.10), Inches(1.95), Inches(0.32), "attacks intercepted", size=11.5, colour=INK, bold=True)
    add_rule(slide, Inches(8.98), Inches(4.64), Inches(3.08), LINE, height=1)
    add_textbox(slide, Inches(8.98), Inches(4.88), Inches(1.32), Inches(0.45), "318/318", size=24, colour=TEAL, bold=True)
    add_textbox(slide, Inches(10.34), Inches(4.94), Inches(1.95), Inches(0.32), "faithful controls accepted", size=11.5, colour=INK, bold=True)
    add_panel(slide, Inches(8.98), Inches(5.62), Inches(3.08), Inches(0.62), fill=PALE, line=PALE, radius=False)
    add_textbox(slide, Inches(9.14), Inches(5.78), Inches(2.76), Inches(0.28),
                "Results apply to the closed, versioned contract.", size=10.5, colour=MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.90), Inches(6.76), Inches(11.55), Inches(0.22),
                "Every detected failure activated fallback. The policy changed what reached the analyst without erasing the raw failure.",
                size=12, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def add_screenshot_panel(slide, x, y, w, h, title: str, message: str, path: Path, crop: tuple[float, float, float, float], colour) -> None:
    add_panel(slide, x, y, w, h, fill=WHITE, line=colour)
    add_textbox(slide, x + Inches(0.16), y + Inches(0.16), w - Inches(0.32), Inches(0.32), title,
                size=14, colour=colour, bold=True)
    left, top, right, bottom = crop
    add_picture_crop(slide, path, x + Inches(0.15), y + Inches(0.62), w - Inches(0.30), h - Inches(1.35),
                     left=left, top=top, right=right, bottom=bottom)
    add_textbox(slide, x + Inches(0.16), y + h - Inches(0.58), w - Inches(0.32), Inches(0.34), message,
                size=9.5, colour=MUTED, align=PP_ALIGN.CENTER)


def build_workbench_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "System demonstration")
    add_badge(slide, Inches(0.68), Inches(1.12), "S0 OPERATIONAL VIEW", fill=TEAL, w=1.65)
    add_screenshot_panel(
        slide, Inches(0.68), Inches(1.55), Inches(4.05), Inches(4.86),
        "1  Model evidence", "Reason codes and deterministic brief remain visible.",
        WORKBENCH, (0.17, 0.28, 0.27, 0.23), BLUE,
    )
    add_screenshot_panel(
        slide, Inches(4.88), Inches(1.55), Inches(4.05), Inches(4.86),
        "2  Validation and fallback", "Rejected candidates cannot become the official brief.",
        GUARDRAIL_FAILURE, (0.41, 0.28, 0.01, 0.01), RUST,
    )
    add_screenshot_panel(
        slide, Inches(9.08), Inches(1.55), Inches(3.57), Inches(4.86),
        "3  Analyst workflow", "Actions and notes are stored in a separate workflow plane.",
        WORKBENCH, (0.73, 0.24, 0.01, 0.08), TEAL,
    )
    add_panel(slide, Inches(1.70), Inches(6.58), Inches(9.95), Inches(0.42), fill=NAVY, line=NAVY, radius=False)
    add_textbox(slide, Inches(1.95), Inches(6.68), Inches(9.45), Inches(0.20),
                "Analyst work is writable. Recorded research evidence is not.",
                size=13, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def add_example_brief(slide, x, title: str, body: str, *, colour, soft, metric: str, metric_label: str) -> None:
    add_panel(slide, x, Inches(1.82), Inches(5.66), Inches(3.66), fill=WHITE, line=colour)
    add_badge(slide, x + Inches(0.24), Inches(2.08), title.upper(), fill=colour, w=1.70)
    add_panel(slide, x + Inches(0.24), Inches(2.68), Inches(5.18), Inches(2.25), fill=soft, line=soft, radius=False)
    add_textbox(slide, x + Inches(0.46), Inches(2.91), Inches(4.72), Inches(1.82), body,
                size=11.5, colour=INK)
    add_textbox(slide, x + Inches(0.26), Inches(5.02), Inches(1.60), Inches(0.34), metric,
                size=22, colour=colour, bold=True)
    add_textbox(slide, x + Inches(1.82), Inches(5.09), Inches(3.45), Inches(0.28), metric_label,
                size=11.5, colour=MUTED)


def build_llm_tradeoff_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "S0 semantic result")
    add_claim_banner(slide, "Shorter articulation came with less visible evidence.", y=1.08, fill=NAVY)
    deterministic = (
        "This above-threshold synthetic alert has Medium relative review priority.\n\n"
        "1. Terminal distance from customer home raises risk.\n"
        "2. Amount versus customer 30-day average raises risk.\n"
        "3. Transaction amount raises risk."
    )
    llm = "All supplied signals raise risk, led by Terminal distance from customer home."
    add_example_brief(slide, Inches(0.72), "Deterministic brief", deterministic,
                      colour=BLUE, soft=BLUE_SOFT, metric="39.09", metric_label="mean words | 3 evidence items")
    add_example_brief(slide, Inches(6.95), "Guarded LLM brief", llm,
                      colour=TEAL, soft=TEAL_SOFT, metric="12.00", metric_label="mean words | 1 evidence item")
    add_arrow(slide, 6.43, 3.20, w=0.40, text="↔", colour=MUTED, size=24)
    add_panel(slide, Inches(1.48), Inches(5.84), Inches(10.38), Inches(0.70), fill=AMBER_SOFT, line=AMBER_SOFT, radius=False)
    add_textbox(slide, Inches(1.74), Inches(6.02), Inches(9.86), Inches(0.30),
                "All 23 accepted candidates selected the shorter option. Two detailed attempts corrupted structured fields and were rejected.",
                size=12.2, colour=AMBER, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.95), Inches(6.67), Inches(9.45), Inches(0.22),
                "The LLM changed presentation, not the evidence available to the analyst.",
                size=12, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_human_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Descriptive human pilot")
    add_panel(slide, Inches(0.68), Inches(1.10), Inches(12.0), Inches(0.48), fill=PALE, line=PALE, radius=False)
    add_textbox(slide, Inches(0.92), Inches(1.23), Inches(11.52), Inches(0.22),
                "11 proxy reviewers | 99 case reviews | descriptive findings, not a professional-analyst performance claim",
                size=11.5, colour=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_picture_contain(slide, HUMAN_PREFERENCES, Inches(0.62), Inches(1.82), Inches(7.72), Inches(4.78))
    add_metric_card(slide, Inches(8.62), Inches(1.84), Inches(1.72), Inches(1.56), "4 / 5", "Median clarity", colour=BLUE)
    add_metric_card(slide, Inches(10.55), Inches(1.84), Inches(1.72), Inches(1.56), "7 / 11", "LLM clearest", colour=TEAL)
    add_metric_card(slide, Inches(8.62), Inches(3.62), Inches(1.72), Inches(1.56), "7 / 11", "LLM first pass", colour=TEAL)
    add_metric_card(slide, Inches(10.55), Inches(3.62), Inches(1.72), Inches(1.56), "6 / 11", "Deterministic trust", colour=BLUE)
    add_panel(slide, Inches(8.62), Inches(5.40), Inches(3.65), Inches(0.95), fill=RUST_SOFT, line=RUST_SOFT, radius=False)
    add_textbox(slide, Inches(8.85), Inches(5.61), Inches(3.20), Inches(0.50),
                "Objective comprehension did not improve with the guarded LLM brief.",
                size=11.5, colour=RUST, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.92), Inches(6.72), Inches(11.55), Inches(0.22),
                "Preference, trust, and evidence comprehension were related but not interchangeable.",
                size=12.3, colour=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build_conclusion_slide(slide, number: int, title: str, script: str) -> None:
    add_header(slide, number, title, "Conclusion")
    cards = [
        ("Detector", "Reproducible signed model evidence, with no claimed autoencoder breakthrough.", BLUE, BLUE_SOFT),
        ("Guardrail", "Raw failures remain measurable and detected failures activate fallback.", RUST, RUST_SOFT),
        ("Local LLM", "Optional articulation only. It does not classify or create fraud evidence.", TEAL, TEAL_SOFT),
        ("Workbench", "Operational review state stays separate from immutable experiment artifacts.", AMBER, AMBER_SOFT),
    ]
    for index, (heading, body, colour, soft) in enumerate(cards):
        row, col = divmod(index, 2)
        x = Inches(0.72 + col * 6.20)
        y = Inches(1.28 + row * 1.72)
        add_panel(slide, x, y, Inches(5.68), Inches(1.42), fill=soft, line=soft)
        add_textbox(slide, x + Inches(0.24), y + Inches(0.22), Inches(1.42), Inches(0.30), heading,
                    size=15, colour=colour, bold=True)
        add_textbox(slide, x + Inches(1.70), y + Inches(0.20), Inches(3.72), Inches(0.78), body,
                    size=12.5, colour=INK)
    add_panel(slide, Inches(0.72), Inches(4.88), Inches(12.0), Inches(0.72), fill=PALE, line=LINE, radius=False)
    add_textbox(slide, Inches(0.96), Inches(5.04), Inches(11.52), Inches(0.36),
                "Limits: anonymous real features, synthetic readable data, one manual auditor, one local 8B model, and a small proxy pilot.",
                size=11.5, colour=MUTED, align=PP_ALIGN.CENTER)
    add_panel(slide, Inches(0.72), Inches(5.86), Inches(12.0), Inches(0.78), fill=NAVY, line=NAVY, radius=False)
    add_textbox(slide, Inches(1.02), Inches(6.08), Inches(11.40), Inches(0.34),
                "Generated fraud-alert text remains untrusted until deterministic checks confirm that it matches the approved evidence.",
                size=15, colour=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.95), Inches(6.76), Inches(11.45), Inches(0.20),
                "Next: independent audit replication, another local model, and real transaction data with interpretable fields.",
                size=10.5, colour=MUTED, align=PP_ALIGN.CENTER)
    add_footer(slide)
    add_notes(slide, script)


def build() -> None:
    data = parse_content()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Fraud Alert Explanations That Fail Closed"
    prs.core_properties.subject = "Sunway University Capstone Project 2 viva presentation"
    prs.core_properties.author = "Ng Yi Zhen"
    blank = prs.slide_layouts[6]

    builders = {
        2: build_problem_slide,
        3: build_gap_slide,
        4: build_roles_slide,
        5: build_method_slide,
        6: build_detector_slide,
        7: build_innovation_slide,
        8: build_guardrail_results_slide,
        9: build_workbench_slide,
        10: build_llm_tradeoff_slide,
        11: build_human_slide,
        12: build_conclusion_slide,
    }

    for slide_data in data:
        number = int(slide_data["number"])
        title = str(slide_data["title"])
        script = str(slide_data["script"])
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = PAPER
        if number == 1:
            build_title_slide(slide, script)
        else:
            builders[number](slide, number, title, script)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    SCRIPT_OUTPUT.write_text(CONTENT.read_text(encoding="utf-8"), encoding="utf-8")
    print(f"Built {OUTPUT}")
    print(f"Updated {SCRIPT_OUTPUT}")


if __name__ == "__main__":
    build()
